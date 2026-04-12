"""PowerPoint<->PDF service. PPTX->PDF uses LibreOffice. PDF->PPTX uses PyMuPDF+python-pptx."""
import fitz
import subprocess
import shutil
import logging
from pathlib import Path
from typing import List
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def pptx_to_pdf(input_path: Path, output_path: Path) -> Path:
      lo = _find_libreoffice()
      if not lo:
                raise RuntimeError("LibreOffice required. Install: sudo apt-get install -y libreoffice")
            res = subprocess.run(
                      [lo, "--headless", "--convert-to", "pdf", "--outdir", str(output_path.parent), str(input_path)],
                      capture_output=True, text=True, timeout=120
            )
    if res.returncode != 0:
              raise RuntimeError(f"LibreOffice error: {res.stderr or res.stdout}")
          converted = output_path.parent / f"{input_path.stem}.pdf"
    if not converted.exists():
              raise RuntimeError("LibreOffice did not produce output file.")
          if converted != output_path:
                    converted.rename(output_path)
                return output_path


def pdf_to_pptx(input_path: Path, output_path: Path, dpi: int = 150) -> Path:
      zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    tmp_images: List[Path] = []
    with fitz.open(str(input_path)) as doc:
              for i, page in enumerate(doc):
                            pix = page.get_pixmap(matrix=mat, alpha=False)
                            img = input_path.parent / f"_tmp_{i}.png"
                            pix.save(str(img))
                            tmp_images.append(img)
                            slide = prs.slides.add_slide(blank)
                            slide.shapes.add_picture(str(img), 0, 0, prs.slide_width, prs.slide_height)
                    prs.save(str(output_path))
    for img in tmp_images:
              try:
                            img.unlink()
except Exception:
            pass
    return output_path


def _find_libreoffice():
      for c in ["libreoffice", "soffice", "/usr/bin/libreoffice", "/usr/bin/soffice",
                              "/usr/lib/libreoffice/program/soffice"]:
                                        if shutil.which(c):
                                                      return c
                                              return None
